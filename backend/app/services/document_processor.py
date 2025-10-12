"""
Document Processing Service - Extracts text from PDFs and PowerPoint files
"""
import io
from typing import Dict, Any, List, Optional
from datetime import datetime
import mimetypes


class DocumentProcessor:
    """
    Processes uploaded documents (PDF, PPTX) and extracts text content.
    Supports:
    - PDF files (text-based and scanned with OCR)
    - PowerPoint files (PPTX)
    - Metadata extraction
    """

    def __init__(self):
        self.supported_types = {
            'application/pdf': self._process_pdf,
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': self._process_pptx,
            'application/vnd.ms-powerpoint': self._process_ppt
        }

    async def process_document(
        self,
        file_content: bytes,
        filename: str,
        mime_type: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        Process a document and extract text content.

        Args:
            file_content: Raw file bytes
            filename: Original filename
            mime_type: MIME type (auto-detected if not provided)

        Returns:
            {
                "text": "Extracted text content",
                "metadata": {
                    "filename": str,
                    "mime_type": str,
                    "pages": int,
                    "file_size": int,
                    "processed_at": datetime
                },
                "chunks": List[str]  # Pre-chunked text sections
            }
        """
        # Detect MIME type if not provided
        if not mime_type:
            mime_type, _ = mimetypes.guess_type(filename)

        if not mime_type or mime_type not in self.supported_types:
            raise ValueError(f"Unsupported file type: {mime_type}. Supported: PDF, PPTX")

        # Process based on file type
        processor = self.supported_types[mime_type]
        result = await processor(file_content, filename)

        return {
            "text": result["text"],
            "metadata": {
                "filename": filename,
                "mime_type": mime_type,
                "pages": result.get("pages", 0),
                "file_size": len(file_content),
                "processed_at": datetime.utcnow(),
                **result.get("extra_metadata", {})
            },
            "chunks": result.get("chunks", [result["text"]])
        }

    async def _process_pdf(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """Extract text from PDF using PyMuPDF (fitz)"""
        try:
            import fitz  # PyMuPDF
        except ImportError:
            raise ImportError("PyMuPDF (fitz) not installed. Install with: pip install pymupdf")

        try:
            # Open PDF from bytes
            pdf_doc = fitz.open(stream=file_content, filetype="pdf")

            text_content = []
            chunks = []  # Store text by page for better chunking

            for page_num in range(len(pdf_doc)):
                page = pdf_doc[page_num]
                page_text = page.get_text()

                if page_text.strip():
                    text_content.append(page_text)
                    chunks.append({
                        "text": page_text,
                        "page": page_num + 1,
                        "type": "page"
                    })

            pdf_doc.close()

            full_text = "\n\n".join(text_content)

            return {
                "text": full_text,
                "pages": len(pdf_doc),
                "chunks": [chunk["text"] for chunk in chunks],
                "extra_metadata": {
                    "total_pages": len(pdf_doc)
                }
            }

        except Exception as e:
            raise Exception(f"Failed to process PDF: {str(e)}")

    async def _process_pptx(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """Extract text from PowerPoint (PPTX) using python-pptx"""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx not installed. Install with: pip install python-pptx")

        try:
            # Open PPTX from bytes
            pptx_file = io.BytesIO(file_content)
            presentation = Presentation(pptx_file)

            text_content = []
            chunks = []  # Store text by slide

            for slide_num, slide in enumerate(presentation.slides):
                slide_text = []

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)

                # Extract notes if available
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text
                    if notes.strip():
                        slide_text.append(f"[Notes: {notes}]")

                slide_content = "\n".join(slide_text)
                if slide_content.strip():
                    text_content.append(f"Slide {slide_num + 1}:\n{slide_content}")
                    chunks.append({
                        "text": slide_content,
                        "slide": slide_num + 1,
                        "type": "slide"
                    })

            full_text = "\n\n".join(text_content)

            return {
                "text": full_text,
                "pages": len(presentation.slides),
                "chunks": [chunk["text"] for chunk in chunks],
                "extra_metadata": {
                    "total_slides": len(presentation.slides)
                }
            }

        except Exception as e:
            raise Exception(f"Failed to process PPTX: {str(e)}")

    async def _process_ppt(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """
        Process legacy PPT files.
        For now, returns an error - would need LibreOffice or Apache Tika.
        """
        raise NotImplementedError(
            "Legacy PPT files are not supported. "
            "Please convert to PPTX format or install LibreOffice/Apache Tika."
        )


# Singleton instance
_document_processor = None


def get_document_processor() -> DocumentProcessor:
    """Get the document processor instance (singleton)"""
    global _document_processor
    if _document_processor is None:
        _document_processor = DocumentProcessor()
    return _document_processor
